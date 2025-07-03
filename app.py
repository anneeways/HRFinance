import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import groq
import json
from datetime import datetime
import io
import base64
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Page configuration
st.set_page_config(
    page_title="AI-Powered HR Kostenvergleich",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize Groq client
@st.cache_resource
def init_groq():
    try:
        api_key = st.secrets.get("GROQ_API_KEY") or st.session_state.get("groq_api_key")
        if not api_key:
            return None
        return groq.Groq(api_key=api_key)
    except Exception as e:
        st.error(f"KI-Initialisierung Fehler: {e}")
        return None

# Industry templates
INDUSTRY_TEMPLATES = {
    "Tech": {
        "hire_salary": 85000,
        "vacancy_months": 4,
        "social_percent": 22,
        "benefits_percent": 12,
        "prod_loss_percent": 50,
        "berater_percent": 30,
        "interview_hours": 15,
        "interview_rate": 80,
        "training_cost": 2000,
    },
    "Healthcare": {
        "hire_salary": 65000,
        "vacancy_months": 3,
        "social_percent": 24,
        "benefits_percent": 15,
        "prod_loss_percent": 40,
        "berater_percent": 20,
        "interview_hours": 10,
        "interview_rate": 70,
        "training_cost": 1500,
    },
    "Retail": {
        "hire_salary": 35000,
        "vacancy_months": 2,
        "social_percent": 20,
        "benefits_percent": 8,
        "prod_loss_percent": 30,
        "berater_percent": 15,
        "interview_hours": 8,
        "interview_rate": 50,
        "training_cost": 500,
    },
    "Finance": {
        "hire_salary": 75000,
        "vacancy_months": 5,
        "social_percent": 23,
        "benefits_percent": 18,
        "prod_loss_percent": 45,
        "berater_percent": 25,
        "interview_hours": 12,
        "interview_rate": 90,
        "training_cost": 2500,
    }
}

def create_pdf_report(results, context_data, ai_insights=None, ai_scenarios=None):
    """Generate a professional PDF report"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.darkblue,
        alignment=1  # Center
    )
    story.append(Paragraph("HR Kostenvergleich Analyse", title_style))
    story.append(Spacer(1, 20))
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", styles['Heading2']))
    
    if results['total_hire'] > results['total_salary_increase']:
        recommendation = f"Gehaltserhöhung ist günstiger um {results['total_hire'] - results['total_salary_increase']:,.0f} €"
        color = colors.green
    else:
        recommendation = f"Neubesetzung ist günstiger um {results['total_salary_increase'] - results['total_hire']:,.0f} €"
        color = colors.blue
    
    summary_style = ParagraphStyle('Summary', parent=styles['Normal'], backColor=color, textColor=colors.white, leftIndent=10, rightIndent=10, spaceBefore=10, spaceAfter=10)
    story.append(Paragraph(f"<b>Empfehlung: {recommendation}</b>", summary_style))
    story.append(Spacer(1, 20))
    
    # Key Metrics Table
    story.append(Paragraph("Kostenübersicht", styles['Heading2']))
    
    data = [
        ['Kostenart', 'Betrag (€)', 'Anteil (%)'],
        ['Neubesetzung Gesamt', f"{results['total_hire']:,.0f}", '100%'],
        ['- Recruiting', f"{results['recruiting']['sum']:,.0f}", f"{results['recruiting']['sum']/results['total_hire']*100:.1f}%"],
        ['- Vakanz', f"{results['vacancy']['sum']:,.0f}", f"{results['vacancy']['sum']/results['total_hire']*100:.1f}%"],
        ['- Onboarding', f"{results['onboarding']['sum']:,.0f}", f"{results['onboarding']['sum']/results['total_hire']*100:.1f}%"],
        ['- Produktivitätsverlust', f"{results['productivity']['sum']:,.0f}", f"{results['productivity']['sum']/results['total_hire']*100:.1f}%"],
        ['- Weitere Kosten', f"{results['other']['sum']:,.0f}", f"{results['other']['sum']/results['total_hire']*100:.1f}%"],
        ['- Fixkosten', f"{results['fixed']['sum']:,.0f}", f"{results['fixed']['sum']/results['total_hire']*100:.1f}%"],
        ['', '', ''],
        ['Gehaltserhöhung Gesamt', f"{results['total_salary_increase']:,.0f}", '100%'],
    ]
    
    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -2), colors.beige),
        ('BACKGROUND', (0, -1), (-1, -1), colors.lightblue),
        ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(table)
    story.append(Spacer(1, 20))
    
    # Parameters
    story.append(Paragraph("Parameter", styles['Heading2']))
    param_data = [
        ['Parameter', 'Wert'],
        ['Jahresgehalt (Neubesetzung)', f"{context_data.get('hire_salary', 0):,.0f} €"],
        ['Vakanzdauer', f"{context_data.get('vacancy_months', 0)} Monate"],
        ['Branche', context_data.get('industry', 'General')],
        ['Produktivitätsverlust', f"{context_data.get('prod_loss_percent', 0)}%"],
        ['Analyse-Datum', datetime.now().strftime('%d.%m.%Y %H:%M')]
    ]
    
    param_table = Table(param_data)
    param_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(param_table)
    
    # AI Insights if available
    if ai_insights:
        story.append(Spacer(1, 20))
        story.append(Paragraph("KI-Strategieanalyse", styles['Heading2']))
        story.append(Paragraph(ai_insights, styles['Normal']))
    
    if ai_scenarios:
        story.append(Spacer(1, 20))
        story.append(Paragraph("KI-Szenarien", styles['Heading2']))
        story.append(Paragraph(ai_scenarios, styles['Normal']))
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def get_ai_insights(groq_client, calculation_data, context_data):

def create_excel_report(results, context_data, ai_insights=None):
    """Generate a comprehensive Excel report"""
    buffer = io.BytesIO()
    
    with xlsxwriter.Workbook(buffer) as workbook:
        # Formats
        title_format = workbook.add_format({'bold': True, 'font_size': 16, 'align': 'center'})
        header_format = workbook.add_format({'bold': True, 'bg_color': '#4472C4', 'font_color': 'white'})
        currency_format = workbook.add_format({'num_format': '#,##0 €'})
        percent_format = workbook.add_format({'num_format': '0.0%'})
        
        # Summary Sheet
        ws1 = workbook.add_worksheet('Zusammenfassung')
        ws1.write('A1', 'HR Kostenvergleich Analyse', title_format)
        ws1.write('A3', 'Erstelldatum:', header_format)
        ws1.write('B3', datetime.now().strftime('%d.%m.%Y %H:%M'))
        
        # Key metrics
        ws1.write('A5', 'Kostenkategorie', header_format)
        ws1.write('B5', 'Betrag (€)', header_format)
        ws1.write('C5', 'Anteil (%)', header_format)
        
        row = 6
        ws1.write(row, 0, 'Neubesetzung Gesamt')
        ws1.write(row, 1, results['total_hire'], currency_format)
        ws1.write(row, 2, 1.0, percent_format)
        
        categories = [
            ('Recruiting', results['recruiting']['sum']),
            ('Vakanz', results['vacancy']['sum']),
            ('Onboarding', results['onboarding']['sum']),
            ('Produktivitätsverlust', results['productivity']['sum']),
            ('Weitere Kosten', results['other']['sum']),
            ('Fixkosten', results['fixed']['sum'])
        ]
        
        for cat, amount in categories:
            row += 1
            ws1.write(row, 0, f"  - {cat}")
            ws1.write(row, 1, amount, currency_format)
            ws1.write(row, 2, amount/results['total_hire'], percent_format)
        
        row += 2
        ws1.write(row, 0, 'Gehaltserhöhung Gesamt')
        ws1.write(row, 1, results['total_salary_increase'], currency_format)
        
        # Detailed breakdown sheet
        ws2 = workbook.add_worksheet('Detailkosten')
        ws2.write('A1', 'Detaillierte Kostenaufschlüsselung', title_format)
        
        # Recruiting details
        row = 3
        ws2.write(row, 0, 'Recruiting-Kosten', header_format)
        row += 1
        for item, cost in results['recruiting']['costs'].items():
            ws2.write(row, 0, item)
            ws2.write(row, 1, cost, currency_format)
            row += 1
        
        # Parameters sheet
        ws3 = workbook.add_worksheet('Parameter')
        ws3.write('A1', 'Eingabeparameter', title_format)
        
        params = [
            ('Jahresgehalt (Neubesetzung)', f"{context_data.get('hire_salary', 0):,.0f} €"),
            ('Vakanzdauer', f"{context_data.get('vacancy_months', 0)} Monate"),
            ('Branche', context_data.get('industry', 'General')),
            ('Produktivitätsverlust', f"{context_data.get('prod_loss_percent', 0)}%")
        ]
        
        row = 3
        for param, value in params:
            ws3.write(row, 0, param)
            ws3.write(row, 1, value)
            row += 1
        
        # Auto-adjust column widths
        for ws in [ws1, ws2, ws3]:
            ws.set_column('A:A', 25)
            ws.set_column('B:B', 15)
            ws.set_column('C:C', 12)
    
    buffer.seek(0)
    return buffer

def create_powerpoint_report(results, context_data, ai_insights=None):
    """Generate a PowerPoint presentation"""
    prs = Presentation()
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "HR Kostenvergleich Analyse"
    subtitle.text = f"Erstellt am {datetime.now().strftime('%d.%m.%Y')}\nBranche: {context_data.get('industry', 'General')}"
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    
    if results['total_hire'] > results['total_salary_increase']:
        recommendation = f"Gehaltserhöhung ist günstiger\nErsparnis: {results['total_hire'] - results['total_salary_increase']:,.0f} €"
    else:
        recommendation = f"Neubesetzung ist günstiger\nErsparnis: {results['total_salary_increase'] - results['total_hire']:,.0f} €"
    
    content = slide2.placeholders[1]
    content.text = f"""Empfehlung: {recommendation}

Neubesetzung Gesamtkosten: {results['total_hire']:,.0f} €
Gehaltserhöhung Kosten: {results['total_salary_increase']:,.0f} €

Größte Kostenfaktoren:
• Fixkosten: {results['fixed']['sum']:,.0f} €
• Produktivitätsverlust: {results['productivity']['sum']:,.0f} €
• Recruiting: {results['recruiting']['sum']:,.0f} €"""
    
    # Slide 3: Cost Breakdown
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Kostenaufschlüsselung Neubesetzung"
    
    breakdown_text = f"""Recruiting: {results['recruiting']['sum']:,.0f} € ({results['recruiting']['sum']/results['total_hire']*100:.1f}%)
Vakanz: {results['vacancy']['sum']:,.0f} € ({results['vacancy']['sum']/results['total_hire']*100:.1f}%)
Onboarding: {results['onboarding']['sum']:,.0f} € ({results['onboarding']['sum']/results['total_hire']*100:.1f}%)
Produktivitätsverlust: {results['productivity']['sum']:,.0f} € ({results['productivity']['sum']/results['total_hire']*100:.1f}%)
Weitere Kosten: {results['other']['sum']:,.0f} € ({results['other']['sum']/results['total_hire']*100:.1f}%)
Fixkosten: {results['fixed']['sum']:,.0f} € ({results['fixed']['sum']/results['total_hire']*100:.1f}%)

Gesamtsumme: {results['total_hire']:,.0f} €"""
    
    slide3.placeholders[1].text = breakdown_text
    
    # Slide 4: AI Insights (if available)
    if ai_insights:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "KI-Strategieanalyse"
        slide4.placeholders[1].text = ai_insights[:1000] + "..." if len(ai_insights) > 1000 else ai_insights
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def create_word_report(results, context_data, ai_insights=None, ai_scenarios=None):
    """Generate a Word document report"""
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = Document()
    
    # Title
    title = doc.add_heading('HR Kostenvergleich Analyse', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Executive Summary
    doc.add_heading('Executive Summary', level=1)
    
    if results['total_hire'] > results['total_salary_increase']:
        recommendation = f"Gehaltserhöhung ist günstiger um {results['total_hire'] - results['total_salary_increase']:,.0f} €"
    else:
        recommendation = f"Neubesetzung ist günstiger um {results['total_salary_increase'] - results['total_hire']:,.0f} €"
    
    summary_para = doc.add_paragraph()
    summary_para.add_run(f"Empfehlung: {recommendation}").bold = True
    
    # Cost table
    doc.add_heading('Kostenübersicht', level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Kostenart'
    hdr_cells[1].text = 'Betrag (€)'
    hdr_cells[2].text = 'Anteil (%)'
    
    cost_items = [
        ('Neubesetzung Gesamt', results['total_hire'], 100),
        ('- Recruiting', results['recruiting']['sum'], results['recruiting']['sum']/results['total_hire']*100),
        ('- Vakanz', results['vacancy']['sum'], results['vacancy']['sum']/results['total_hire']*100),
        ('- Onboarding', results['onboarding']['sum'], results['onboarding']['sum']/results['total_hire']*100),
        ('- Produktivitätsverlust', results['productivity']['sum'], results['productivity']['sum']/results['total_hire']*100),
        ('- Weitere Kosten', results['other']['sum'], results['other']['sum']/results['total_hire']*100),
        ('- Fixkosten', results['fixed']['sum'], results['fixed']['sum']/results['total_hire']*100),
        ('Gehaltserhöhung Gesamt', results['total_salary_increase'], 100),
    ]
    
    for item, amount, percent in cost_items:
        row_cells = table.add_row().cells
        row_cells[0].text = item
        row_cells[1].text = f"{amount:,.0f}"
        row_cells[2].text = f"{percent:.1f}%"
    
    # Parameters
    doc.add_heading('Parameter', level=1)
    param_para = doc.add_paragraph()
    param_para.add_run(f"Jahresgehalt (Neubesetzung): {context_data.get('hire_salary', 0):,.0f} €\n")
    param_para.add_run(f"Vakanzdauer: {context_data.get('vacancy_months', 0)} Monate\n")
    param_para.add_run(f"Branche: {context_data.get('industry', 'General')}\n")
    param_para.add_run(f"Produktivitätsverlust: {context_data.get('prod_loss_percent', 0)}%\n")
    param_para.add_run(f"Analyse-Datum: {datetime.now().strftime('%d.%m.%Y %H:%M')}")
    
    # AI sections
    if ai_insights:
        doc.add_heading('KI-Strategieanalyse', level=1)
        doc.add_paragraph(ai_insights)
    
    if ai_scenarios:
        doc.add_heading('KI-Szenarien', level=1)
        doc.add_paragraph(ai_scenarios)
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer
    """Get AI-powered insights using Groq"""
    if not groq_client:
        return None
    
    try:
        prompt = f"""
        Als HR-Experte analysiere bitte diese Kostenvergleichsdaten und gib strategische Empfehlungen:

        KOSTENDATEN:
        - Neubesetzung Gesamtkosten: {calculation_data['total_hire']:,.0f} €
        - Gehaltserhöhung Gesamtkosten: {calculation_data['total_salary_increase']:,.0f} €
        - Jahresgehalt: {context_data['hire_salary']:,.0f} €
        - Branche: {context_data.get('industry', 'Unbekannt')}
        - Vakanzdauer: {context_data['vacancy_months']} Monate
        - Produktivitätsverlust: {context_data['prod_loss_percent']}%

        KOSTENAUFSCHLÜSSELUNG:
        - Recruiting: {calculation_data['recruiting']['sum']:,.0f} €
        - Vakanz: {calculation_data['vacancy']['sum']:,.0f} €
        - Onboarding: {calculation_data['onboarding']['sum']:,.0f} €
        - Produktivitätsverlust: {calculation_data['productivity']['sum']:,.0f} €
        - Weitere Kosten: {calculation_data['other']['sum']:,.0f} €

        Bitte analysiere und gib zurück:
        1. Strategische Empfehlung (Neubesetzung vs Gehaltserhöhung)
        2. Top 3 Kostentreiber identifizieren
        3. Konkrete Optimierungsvorschläge
        4. Risikobewertung für beide Optionen
        5. Langzeit-Perspektive (3-5 Jahre)

        Antworte auf Deutsch, präzise und geschäftsorientiert.
        """

        chat_completion = groq_client.chat.completions.create(
            messages=[{
                "role": "system",
                "content": "Du bist ein erfahrener HR-Strategieberater mit 15+ Jahren Erfahrung in Personalkosten-Optimierung."
            }, {
                "role": "user", 
                "content": prompt
            }],
            model="llama-3.1-8b-instant",
            temperature=0.3,
            max_tokens=1000
        )
        
        return chat_completion.choices[0].message.content
    except Exception as e:
        st.error(f"KI-Analyse Fehler: {e}")
        return None

def get_ai_scenarios(groq_client, calculation_data):
    """Generate AI-powered what-if scenarios"""
    if not groq_client:
        return None
    
    try:
        prompt = f"""
        Erstelle 3 realistische What-If-Szenarien für diesen HR-Kostenvergleich:

        BASISDATEN:
        - Neubesetzung: {calculation_data['total_hire']:,.0f} €
        - Gehaltserhöhung: {calculation_data['total_salary_increase']:,.0f} €

        Erstelle Szenarien für:
        1. Best-Case (optimistische Annahmen)
        2. Worst-Case (pessimistische Annahmen)  
        3. Economic Downturn (Wirtschaftskrise)

        Für jedes Szenario gib an:
        - Kurze Beschreibung der Annahmen
        - Geschätzte Kostenveränderung in % 
        - Empfehlung für dieses Szenario

        Format als strukturierten Text, nicht als JSON.
        """

        chat_completion = groq_client.chat.completions.create(
            messages=[{
                "role": "user", 
                "content": prompt
            }],
            model="llama-3.1-8b-instant",
            temperature=0.5,
            max_tokens=800
        )
        
        return chat_completion.choices[0].message.content
    except Exception as e:
        st.error(f"Szenario-Generierung Fehler: {e}")
        return None

def load_template(template_name):
    """Load industry template into session state"""
    template = INDUSTRY_TEMPLATES[template_name]
    for key, value in template.items():
        st.session_state[key] = value
    st.session_state['industry'] = template_name

def reset_to_defaults():
    """Reset all values to defaults"""
    defaults = {
        # Basic assumptions
        "hire_salary": 60000,
        "vacancy_months": 3,
        "social_percent": 22,
        "benefits_percent": 8,
        "prod_loss_percent": 40,
        "industry": "General",
        
        # Recruiting costs
        "anzeigen_qty": 2,
        "anzeigen_price": 800,
        "berater_percent": 25,
        "interview_hours": 12,
        "interview_rate": 70,
        "assessment_qty": 1,
        "assessment_price": 1500,
        "reise_qty": 2,
        "reise_price": 300,
        "background_qty": 1,
        "background_price": 200,
        
        # Vacancy costs
        "produkt_price": 6000,
        "ueberstunden_qty": 30,
        "ueberstunden_price": 50,
        "extern_qty": 20,
        "extern_price": 400,
        "gehalt_price": 6000,
        
        # Onboarding costs
        "hr_hours": 10,
        "hr_rate": 50,
        "kollegen_hours": 15,
        "kollegen_rate": 60,
        "training_cost": 1000,
        "it_cost": 1200,
        "mentor_hours": 6,
        "mentor_rate": 60,
        
        # Other costs
        "fehler_cost": 1400,
        "knowhow_cost": 2000,
        "kunden_cost": 2500,
        "team_cost": 2000,
        
        # Salary increase
        "current_salary": 60000,
        "increase_percent": 8,
        "social_increase_percent": 22,
        "benefits_increase_percent": 8,
    }
    
    for key, value in defaults.items():
        st.session_state[key] = value

def initialize_session_state():
    """Initialize session state with default values"""
    if 'initialized' not in st.session_state:
        reset_to_defaults()
        st.session_state.initialized = True

def calculate_costs():
    """Calculate all costs and return results"""
    # Get values from session state
    hire_salary = st.session_state.get('hire_salary', 60000)
    vacancy_months = st.session_state.get('vacancy_months', 3)
    social_percent = st.session_state.get('social_percent', 22)
    benefits_percent = st.session_state.get('benefits_percent', 8)
    prod_loss_percent = st.session_state.get('prod_loss_percent', 40)
    
    # Recruiting costs
    recruiting_costs = {
        "Stellenanzeigen": st.session_state.get('anzeigen_qty', 2) * st.session_state.get('anzeigen_price', 800),
        "Personalberater": hire_salary * (st.session_state.get('berater_percent', 25) / 100),
        "Interviews": st.session_state.get('interview_hours', 12) * st.session_state.get('interview_rate', 70),
        "Assessment Center": st.session_state.get('assessment_qty', 1) * st.session_state.get('assessment_price', 1500),
        "Reisekosten": st.session_state.get('reise_qty', 2) * st.session_state.get('reise_price', 300),
        "Background Checks": st.session_state.get('background_qty', 1) * st.session_state.get('background_price', 200),
    }
    recruiting_sum = sum(recruiting_costs.values())
    
    # Vacancy costs
    vacancy_costs = {
        "Entgangene Produktivität": vacancy_months * st.session_state.get('produkt_price', 6000),
        "Überstunden Team": st.session_state.get('ueberstunden_qty', 30) * st.session_state.get('ueberstunden_price', 50),
        "Externe Unterstützung": st.session_state.get('extern_qty', 20) * st.session_state.get('extern_price', 400),
        "Gehaltsersparnis": -(vacancy_months * st.session_state.get('gehalt_price', 6000)),
    }
    vacancy_sum = sum(vacancy_costs.values())
    
    # Onboarding costs
    onboarding_costs = {
        "HR-Aufwand": st.session_state.get('hr_hours', 10) * st.session_state.get('hr_rate', 50),
        "Einarbeitung Kollegen": st.session_state.get('kollegen_hours', 15) * st.session_state.get('kollegen_rate', 60),
        "Schulungen/Training": st.session_state.get('training_cost', 1000),
        "IT-Setup & Equipment": st.session_state.get('it_cost', 1200),
        "Mentor/Buddy-System": st.session_state.get('mentor_hours', 6) * st.session_state.get('mentor_rate', 60),
    }
    onboarding_sum = sum(onboarding_costs.values())
    
    # Productivity loss
    prod_loss_monthly = (hire_salary / 12) * (1 + social_percent / 100) * (prod_loss_percent / 100)
    productivity_sum = prod_loss_monthly * vacancy_months
    
    # Other costs
    other_costs = {
        "Fehlerrate": st.session_state.get('fehler_cost', 1400),
        "Know-how-Verlust": st.session_state.get('knowhow_cost', 2000),
        "Kundenbindung/Umsatz": st.session_state.get('kunden_cost', 2500),
        "Team-Moral": st.session_state.get('team_cost', 2000),
    }
    other_sum = sum(other_costs.values())
    
    # Fixed costs
    social_hire = hire_salary * (social_percent / 100)
    benefits_hire = hire_salary * (benefits_percent / 100)
    fixed_sum = hire_salary + social_hire + benefits_hire
    
    # Total hiring cost
    total_hire = recruiting_sum + vacancy_sum + onboarding_sum + productivity_sum + other_sum + fixed_sum
    
    # Salary increase costs
    current_salary = st.session_state.get('current_salary', 60000)
    increase_percent = st.session_state.get('increase_percent', 8)
    social_increase_percent = st.session_state.get('social_increase_percent', 22)
    benefits_increase_percent = st.session_state.get('benefits_increase_percent', 8)
    
    increase_amount = current_salary * (increase_percent / 100)
    social_increase = increase_amount * (social_increase_percent / 100)
    benefits_increase = increase_amount * (benefits_increase_percent / 100)
    total_salary_increase = increase_amount + social_increase + benefits_increase
    
    return {
        "recruiting": {"costs": recruiting_costs, "sum": recruiting_sum},
        "vacancy": {"costs": vacancy_costs, "sum": vacancy_sum},
        "onboarding": {"costs": onboarding_costs, "sum": onboarding_sum},
        "productivity": {"sum": productivity_sum},
        "other": {"costs": other_costs, "sum": other_sum},
        "fixed": {"sum": fixed_sum},
        "total_hire": total_hire,
        "total_salary_increase": total_salary_increase,
        "salary_breakdown": {
            "increase": increase_amount,
            "social": social_increase,
            "benefits": benefits_increase
        }
    }

def main():
    # Initialize session state
    initialize_session_state()
    
    # Header
    st.title("🤖 AI-Powered HR Kostenvergleich")
    st.markdown("""
    **Intelligenter Kostenvergleich** zwischen Neubesetzung und Gehaltserhöhung mit **KI-gestützten Insights**. 
    Alle Werte sind editierbar und werden in Echtzeit mit AI-Empfehlungen aktualisiert.
    """)
    
    # AI API Key input (if not in secrets)
    groq_client = init_groq()
    if not groq_client:
        with st.expander("🔑 AI Setup", expanded=True):
            st.info("Für AI-Features benötigen Sie einen API Key.")
            api_key = st.text_input("AI API Key", type="password", help="API Key für KI-Funktionen")
            if api_key:
                st.session_state.groq_api_key = api_key
                groq_client = init_groq()
                if groq_client:
                    st.success("✅ AI erfolgreich verbunden!")
                    st.rerun()
    
    # Sidebar for main assumptions
    with st.sidebar:
        st.header("⚙️ Grundannahmen")
        
        # Industry template selector
        col1, col2 = st.columns(2)
        with col1:
            template = st.selectbox("🏭 Branche", [""] + list(INDUSTRY_TEMPLATES.keys()))
            if template and st.button("Vorlage laden"):
                load_template(template)
                st.rerun()
        
        with col2:
            if st.button("🔄 Reset"):
                reset_to_defaults()
                st.rerun()
        
        st.divider()
        
        # Basic assumptions
        st.subheader("Grundparameter")
        st.number_input("Jahresgehalt (Neubesetzung) €", 
                       min_value=20000, max_value=200000, step=1000,
                       key="hire_salary")
        
        st.number_input("Vakanzdauer (Monate)", 
                       min_value=1, max_value=24, step=1,
                       key="vacancy_months")
        
        st.number_input("Sozialabgaben (%)", 
                       min_value=15, max_value=30, step=1,
                       key="social_percent")
        
        st.number_input("Benefits (%)", 
                       min_value=5, max_value=25, step=1,
                       key="benefits_percent")
        
        st.slider("Produktivitätsverlust (%)", 
                 min_value=0, max_value=100, step=5,
                 key="prod_loss_percent",
                 help="Wie viel Prozent der Arbeitsleistung fehlen pro Monat während der Einarbeitung?")
        
        # AI Features toggle
        st.divider()
        st.subheader("🤖 AI Features")
        use_ai_insights = st.checkbox("AI-Insights aktivieren", value=bool(groq_client))
        use_ai_scenarios = st.checkbox("AI-Szenarien generieren", value=bool(groq_client))
    
    # Main content area
    results = calculate_costs()
    
    # Top-level metrics
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("💼 Neubesetzung", f"{results['total_hire']:,.0f} €", 
                 delta=f"{results['total_hire'] - results['total_salary_increase']:+,.0f} €")
    with col2:
        st.metric("💰 Gehaltserhöhung", f"{results['total_salary_increase']:,.0f} €")
    with col3:
        difference = abs(results['total_hire'] - results['total_salary_increase'])
        percentage = (difference / min(results['total_hire'], results['total_salary_increase'])) * 100
        st.metric("💡 Ersparnis", f"{difference:,.0f} €", f"{percentage:.1f}%")
    
    # Recommendation
    if results['total_hire'] > results['total_salary_increase']:
        st.success("🎯 **KI-Empfehlung: Gehaltserhöhung ist günstiger**")
        st.info(f"💰 Sie sparen {difference:,.0f} € ({percentage:.1f}%) mit einer Gehaltserhöhung")
    else:
        st.info("🎯 **KI-Empfehlung: Neubesetzung ist günstiger**")
        st.success(f"💰 Sie sparen {difference:,.0f} € ({percentage:.1f}%) mit einer Neubesetzung")
    
    # AI Insights Section
    if groq_client and use_ai_insights:
        with st.container():
            st.header("🧠 KI-gestützte Strategieanalyse")
            
            if st.button("🚀 AI-Analyse generieren", type="primary"):
                with st.spinner("🤖 KI analysiert Ihre Daten..."):
                    context_data = {
                        'hire_salary': st.session_state.get('hire_salary', 60000),
                        'vacancy_months': st.session_state.get('vacancy_months', 3),
                        'prod_loss_percent': st.session_state.get('prod_loss_percent', 40),
                        'industry': st.session_state.get('industry', 'General')
                    }
                    
                    insights = get_ai_insights(groq_client, results, context_data)
                    
                    if insights:
                        st.success("✅ AI-Analyse abgeschlossen!")
                        st.markdown("### 🎯 Strategische Empfehlungen")
                        st.markdown(insights)
                        
                        # Save insights to session state
                        st.session_state.ai_insights = insights
                        st.session_state.insights_timestamp = datetime.now()
            
            # Display cached insights if available
            if hasattr(st.session_state, 'ai_insights'):
                st.markdown("### 📋 Letzte AI-Analyse")
                st.info(f"Erstellt: {st.session_state.insights_timestamp.strftime('%d.%m.%Y %H:%M')}")
                st.markdown(st.session_state.ai_insights)
    
    # Detailed input sections
    col1, col2 = st.columns([2, 1])
    
    with col1:
        # Detailed cost inputs
        st.header("🏢 Neubesetzung - Detailkosten")
        
        # Recruiting costs
        with st.expander("🧲 Recruiting-Kosten", expanded=False):
            col_a, col_b = st.columns(2)
            with col_a:
                st.number_input("Stellenanzeigen (Anzahl)", min_value=0, key="anzeigen_qty")
                st.number_input("Stellenanzeigen (€ pro Anzeige)", min_value=0, key="anzeigen_price")
                st.number_input("Personalberater (%)", min_value=0, max_value=50, key="berater_percent")
            with col_b:
                st.number_input("Interview-Stunden", min_value=0, key="interview_hours")
                st.number_input("Interview-Stundensatz (€)", min_value=0, key="interview_rate")
                st.number_input("Assessment Center (€)", min_value=0, key="assessment_price")
        
        # Vacancy costs  
        with st.expander("⏳ Vakanz-Kosten", expanded=False):
            col_a, col_b = st.columns(2)
            with col_a:
                st.number_input("Entgangene Produktivität (€/Monat)", min_value=0, key="produkt_price")
                st.number_input("Überstunden (Anzahl)", min_value=0, key="ueberstunden_qty")
                st.number_input("Überstunden (€/Std)", min_value=0, key="ueberstunden_price")
            with col_b:
                st.number_input("Externe Unterstützung (Tage)", min_value=0, key="extern_qty")
                st.number_input("Externe Unterstützung (€/Tag)", min_value=0, key="extern_price")
                st.number_input("Monatliches Gehalt (€)", min_value=0, key="gehalt_price")
        
        # Onboarding costs
        with st.expander("🎓 Onboarding-Kosten", expanded=False):
            col_a, col_b = st.columns(2)
            with col_a:
                st.number_input("HR-Aufwand (Stunden)", min_value=0, key="hr_hours")
                st.number_input("HR-Stundensatz (€)", min_value=0, key="hr_rate")
                st.number_input("Einarbeitung Kollegen (Stunden)", min_value=0, key="kollegen_hours")
                st.number_input("Kollegen-Stundensatz (€)", min_value=0, key="kollegen_rate")
            with col_b:
                st.number_input("Schulungen/Training (€)", min_value=0, key="training_cost")
                st.number_input("IT-Setup & Equipment (€)", min_value=0, key="it_cost")
                st.number_input("Mentor-Stunden", min_value=0, key="mentor_hours")
                st.number_input("Mentor-Stundensatz (€)", min_value=0, key="mentor_rate")
        
        # Other costs
        with st.expander("⚠️ Weitere Kosten", expanded=False):
            col_a, col_b = st.columns(2)
            with col_a:
                st.number_input("Fehlerrate (€)", min_value=0, key="fehler_cost")
                st.number_input("Know-how-Verlust (€)", min_value=0, key="knowhow_cost")
            with col_b:
                st.number_input("Kundenbindung/Umsatzverluste (€)", min_value=0, key="kunden_cost")
                st.number_input("Team-Moral (€)", min_value=0, key="team_cost")
        
        # Salary increase section
        st.header("💰 Alternative: Gehaltserhöhung")
        with st.expander("💶 Gehaltserhöhung Details", expanded=True):
            col_a, col_b = st.columns(2)
            with col_a:
                st.number_input("Aktuelles Jahresgehalt (€)", min_value=0, key="current_salary")
                st.number_input("Erhöhung (%)", min_value=0, max_value=50, key="increase_percent")
            with col_b:
                st.number_input("Sozialabgaben auf Erhöhung (%)", min_value=0, key="social_increase_percent")
                st.number_input("Benefits auf Erhöhung (%)", min_value=0, key="benefits_increase_percent")
    
    with col2:
        # Cost breakdown chart
        st.subheader("📊 Kostenverteilung")
        
        categories = ["Recruiting", "Vakanz", "Onboarding", "Produktivität", "Weitere", "Fixkosten"]
        values = [
            results['recruiting']['sum'],
            results['vacancy']['sum'],
            results['onboarding']['sum'],
            results['productivity']['sum'],
            results['other']['sum'],
            results['fixed']['sum']
        ]
        
        fig = px.pie(
            values=values,
            names=categories,
            title="Neubesetzung - Kostenverteilung",
            color_discrete_sequence=px.colors.qualitative.Set3
        )
        fig.update_traces(textposition='inside', textinfo='percent+label')
        fig.update_layout(height=400)
        st.plotly_chart(fig, use_container_width=True)
        
        # Comparison chart
        st.subheader("⚖️ Direktvergleich")
        comparison_data = {
            "Option": ["Neubesetzung", "Gehaltserhöhung"],
            "Kosten": [results['total_hire'], results['total_salary_increase']]
        }
        
        fig2 = px.bar(
            comparison_data,
            x="Option",
            y="Kosten",
            title="Kostenvergleich",
            color="Kosten",
            color_continuous_scale="RdYlGn_r"
        )
        fig2.update_layout(showlegend=False, height=300)
        st.plotly_chart(fig2, use_container_width=True)
    
    # AI Scenarios Section
    if groq_client and use_ai_scenarios:
        st.header("🔮 KI-generierte What-If-Szenarien")
        
        if st.button("🎲 AI-Szenarien generieren"):
            with st.spinner("🤖 KI erstellt Szenarien..."):
                scenarios = get_ai_scenarios(groq_client, results)
                
                if scenarios:
                    st.success("✅ Szenarien generiert!")
                    st.markdown("### 📈 What-If-Szenarien")
                    st.markdown(scenarios)
                    
                    # Save scenarios to session state
                    st.session_state.ai_scenarios = scenarios
                    st.session_state.scenarios_timestamp = datetime.now()
        
        # Display cached scenarios if available
        if hasattr(st.session_state, 'ai_scenarios'):
            st.markdown("### 📋 Letzte AI-Szenarien")
            st.info(f"Erstellt: {st.session_state.scenarios_timestamp.strftime('%d.%m.%Y %H:%M')}")
            st.markdown(st.session_state.ai_scenarios)
    
    # Detailed breakdown at bottom
    st.header("📋 Detaillierte Kostenaufschlüsselung")
    
    tab1, tab2, tab3, tab4 = st.tabs(["💼 Neubesetzung Details", "💰 Gehaltserhöhung Details", "📄 Export", "🤖 AI History"])
    
    with tab1:
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.subheader("🧲 Recruiting")
            for item, cost in results['recruiting']['costs'].items():
                st.write(f"{item}: {cost:,.0f} €")
            st.write(f"**Summe: {results['recruiting']['sum']:,.0f} €**")
            
            st.subheader("🎓 Onboarding")
            for item, cost in results['onboarding']['costs'].items():
                st.write(f"{item}: {cost:,.0f} €")
            st.write(f"**Summe: {results['onboarding']['sum']:,.0f} €**")
        
        with col2:
            st.subheader("⏳ Vakanz")
            for item, cost in results['vacancy']['costs'].items():
                st.write(f"{item}: {cost:,.0f} €")
            st.write(f"**Summe: {results['vacancy']['sum']:,.0f} €**")
            
            st.subheader("⚠️ Weitere Kosten")
            for item, cost in results['other']['costs'].items():
                st.write(f"{item}: {cost:,.0f} €")
            st.write(f"**Summe: {results['other']['sum']:,.0f} €**")
        
        with col3:
            st.subheader("📉 Produktivitätsverlust")
            st.write(f"Monatlicher Verlust: {results['productivity']['sum']/st.session_state.vacancy_months:,.0f} €")
            st.write(f"**Gesamtverlust: {results['productivity']['sum']:,.0f} €**")
            
            st.subheader("💶 Fixkosten")
            st.write(f"Jahresgehalt: {st.session_state.hire_salary:,.0f} €")
            st.write(f"Sozialabgaben: {st.session_state.hire_salary * st.session_state.social_percent/100:,.0f} €")
            st.write(f"Benefits: {st.session_state.hire_salary * st.session_state.benefits_percent/100:,.0f} €")
            st.write(f"**Summe: {results['fixed']['sum']:,.0f} €**")
    
    with tab2:
        breakdown = results['salary_breakdown']
        st.subheader("💰 Gehaltserhöhung Aufschlüsselung")
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Grunderhöhung", f"{breakdown['increase']:,.0f} €")
            st.metric("Sozialabgaben", f"{breakdown['social']:,.0f} €")
            st.metric("Benefits", f"{breakdown['benefits']:,.0f} €")
        with col2:
            st.metric("**Gesamtkosten**", f"**{results['total_salary_increase']:,.0f} €**")
            
            # Visualization of salary increase breakdown
            fig_salary = px.pie(
                values=[breakdown['increase'], breakdown['social'], breakdown['benefits']],
                names=['Grunderhöhung', 'Sozialabgaben', 'Benefits'],
                title="Gehaltserhöhung Aufschlüsselung"
            )
            st.plotly_chart(fig_salary, use_container_width=True)
    
    with tab3:
        st.subheader("📄 Export & Sharing")
        
        st.markdown("### 🚀 Professional Export Options")
        st.info("Wählen Sie das passende Format für Ihren Anwendungsfall:")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**📊 Für Präsentationen:**")
            if st.button("📑 PowerPoint Export", help="Ideal für Stakeholder-Präsentationen"):
                with st.spinner("Erstelle PowerPoint-Präsentation..."):
                    try:
                        context_data = {
                            'hire_salary': st.session_state.get('hire_salary', 60000),
                            'vacancy_months': st.session_state.get('vacancy_months', 3),
                            'industry': st.session_state.get('industry', 'General'),
                            'prod_loss_percent': st.session_state.get('prod_loss_percent', 40)
                        }
                        
                        ai_insights = getattr(st.session_state, 'ai_insights', None)
                        pptx_buffer = create_powerpoint_report(results, context_data, ai_insights)
                        
                        st.download_button(
                            label="📥 PowerPoint herunterladen",
                            data=pptx_buffer,
                            file_name=f"hr_kostenvergleich_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success("✅ PowerPoint-Präsentation erstellt!")
                    except Exception as e:
                        st.error(f"PowerPoint-Export Fehler: {e}")
            
            st.markdown("**📋 Für Dokumentation:**")
            if st.button("📄 PDF Report", help="Professioneller Bericht für Archivierung"):
                with st.spinner("Erstelle PDF-Report..."):
                    try:
                        context_data = {
                            'hire_salary': st.session_state.get('hire_salary', 60000),
                            'vacancy_months': st.session_state.get('vacancy_months', 3),
                            'industry': st.session_state.get('industry', 'General'),
                            'prod_loss_percent': st.session_state.get('prod_loss_percent', 40)
                        }
                        
                        ai_insights = getattr(st.session_state, 'ai_insights', None)
                        ai_scenarios = getattr(st.session_state, 'ai_scenarios', None)
                        pdf_buffer = create_pdf_report(results, context_data, ai_insights, ai_scenarios)
                        
                        st.download_button(
                            label="📥 PDF herunterladen",
                            data=pdf_buffer,
                            file_name=f"hr_kostenvergleich_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                            mime="application/pdf"
                        )
                        st.success("✅ PDF-Report erstellt!")
                    except Exception as e:
                        st.error(f"PDF-Export Fehler: {e}")
        
        with col2:
            st.markdown("**📊 Für Analyse:**")
            if st.button("📊 Excel Export", help="Detaillierte Daten für weitere Analyse"):
                with st.spinner("Erstelle Excel-Datei..."):
                    try:
                        context_data = {
                            'hire_salary': st.session_state.get('hire_salary', 60000),
                            'vacancy_months': st.session_state.get('vacancy_months', 3),
                            'industry': st.session_state.get('industry', 'General'),
                            'prod_loss_percent': st.session_state.get('prod_loss_percent', 40)
                        }
                        
                        ai_insights = getattr(st.session_state, 'ai_insights', None)
                        excel_buffer = create_excel_report(results, context_data, ai_insights)
                        
                        st.download_button(
                            label="📥 Excel herunterladen",
                            data=excel_buffer,
                            file_name=f"hr_kostenvergleich_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        )
                        st.success("✅ Excel-Datei erstellt!")
                    except Exception as e:
                        st.error(f"Excel-Export Fehler: {e}")
            
            st.markdown("**📝 Für Bearbeitung:**")
            if st.button("📝 Word Document", help="Editierbarer Bericht für weitere Bearbeitung"):
                with st.spinner("Erstelle Word-Dokument..."):
                    try:
                        context_data = {
                            'hire_salary': st.session_state.get('hire_salary', 60000),
                            'vacancy_months': st.session_state.get('vacancy_months', 3),
                            'industry': st.session_state.get('industry', 'General'),
                            'prod_loss_percent': st.session_state.get('prod_loss_percent', 40)
                        }
                        
                        ai_insights = getattr(st.session_state, 'ai_insights', None)
                        ai_scenarios = getattr(st.session_state, 'ai_scenarios', None)
                        word_buffer = create_word_report(results, context_data, ai_insights, ai_scenarios)
                        
                        st.download_button(
                            label="📥 Word herunterladen",
                            data=word_buffer,
                            file_name=f"hr_kostenvergleich_{datetime.now().strftime('%Y%m%d_%H%M')}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )
                        st.success("✅ Word-Dokument erstellt!")
                    except Exception as e:
                        st.error(f"Word-Export Fehler: {e}")
        
        st.divider()
        
        # Quick data preview
        st.markdown("### 📋 Datenvorschau")
        
        # Create export data preview
        export_data = {
            "Parameter": [
                "Jahresgehalt Neubesetzung", "Vakanzdauer", "Branche", "Sozialabgaben", 
                "Benefits", "Produktivitätsverlust", "Analyse-Datum"
            ],
            "Wert": [
                f"{st.session_state.hire_salary:,} €", 
                f"{st.session_state.vacancy_months} Monate",
                st.session_state.get('industry', 'General'),
                f"{st.session_state.social_percent}%", 
                f"{st.session_state.benefits_percent}%", 
                f"{st.session_state.prod_loss_percent}%",
                datetime.now().strftime('%d.%m.%Y %H:%M')
            ]
        }
        
        results_data = {
            "Kostenart": [
                "Neubesetzung Gesamt", "Recruiting", "Vakanz", "Onboarding", 
                "Produktivitätsverlust", "Weitere Kosten", "Fixkosten",
                "Gehaltserhöhung Gesamt", "Empfohlene Option"
            ],
            "Betrag": [
                f"{results['total_hire']:,} €",
                f"{results['recruiting']['sum']:,} €",
                f"{results['vacancy']['sum']:,} €", 
                f"{results['onboarding']['sum']:,} €",
                f"{results['productivity']['sum']:,} €",
                f"{results['other']['sum']:,} €",
                f"{results['fixed']['sum']:,} €",
                f"{results['total_salary_increase']:,} €",
                "Gehaltserhöhung" if results['total_hire'] > results['total_salary_increase'] else "Neubesetzung"
            ]
        }
        
        df_params = pd.DataFrame(export_data)
        df_results = pd.DataFrame(results_data)
        
        col1, col2 = st.columns(2)
        with col1:
            st.write("**📊 Parameter:**")
            st.dataframe(df_params, hide_index=True, use_container_width=True)
        
        with col2:
            st.write("**💰 Ergebnisse:**")
            st.dataframe(df_results, hide_index=True, use_container_width=True)
        
        # Legacy CSV export
        st.divider()
        st.markdown("### 📊 Legacy Export")
        
        # Enhanced export with AI insights
        csv_data = "HR KOSTENVERGLEICH - DETAILBERICHT\n"
        csv_data += "=" * 50 + "\n\n"
        csv_data += f"Erstellt am: {datetime.now().strftime('%d.%m.%Y %H:%M')}\n"
        csv_data += f"Branche: {st.session_state.get('industry', 'General')}\n\n"
        csv_data += "PARAMETER:\n" + df_params.to_csv(index=False) + "\n"
        csv_data += "ERGEBNISSE:\n" + df_results.to_csv(index=False) + "\n"
        
        if hasattr(st.session_state, 'ai_insights'):
            csv_data += "\nKI-ANALYSE:\n" + "=" * 20 + "\n"
            csv_data += st.session_state.ai_insights + "\n"
        
        if hasattr(st.session_state, 'ai_scenarios'):
            csv_data += "\nKI-SZENARIEN:\n" + "=" * 20 + "\n"
            csv_data += st.session_state.ai_scenarios + "\n"
        
        st.download_button(
            label="📥 Als CSV herunterladen (Legacy)",
            data=csv_data,
            file_name=f"hr_kostenvergleich_{datetime.now().strftime('%Y%m%d_%H%M')}.csv",
            mime="text/csv"
        )
    
    with tab4:
        st.subheader("🤖 AI-Analyse History")
        
        if hasattr(st.session_state, 'ai_insights'):
            st.markdown("### 🧠 Letzte Strategieanalyse")
            st.info(f"Erstellt: {st.session_state.insights_timestamp.strftime('%d.%m.%Y %H:%M')}")
            with st.expander("Vollständige Analyse anzeigen"):
                st.markdown(st.session_state.ai_insights)
        
        if hasattr(st.session_state, 'ai_scenarios'):
            st.markdown("### 🔮 Letzte Szenarien")
            st.info(f"Erstellt: {st.session_state.scenarios_timestamp.strftime('%d.%m.%Y %H:%M')}")
            with st.expander("Vollständige Szenarien anzeigen"):
                st.markdown(st.session_state.ai_scenarios)
        
        if not hasattr(st.session_state, 'ai_insights') and not hasattr(st.session_state, 'ai_scenarios'):
            st.info("🤷‍♂️ Noch keine AI-Analysen erstellt. Aktivieren Sie die AI-Features und generieren Sie Insights!")
        
        # Clear AI history
        if st.button("🗑️ AI-History löschen"):
            if hasattr(st.session_state, 'ai_insights'):
                del st.session_state.ai_insights
                del st.session_state.insights_timestamp
            if hasattr(st.session_state, 'ai_scenarios'):
                del st.session_state.ai_scenarios
                del st.session_state.scenarios_timestamp
            st.success("✅ AI-History gelöscht!")
            st.rerun()

    # Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #666; padding: 20px;'>
        <p>🤖 Powered by <strong>Künstliche Intelligenz</strong> | 🐍 Built with <strong>Streamlit</strong> | 💼 HR Intelligence Platform</p>
        <p><small>Alle Berechnungen sind Schätzungen. Konsultieren Sie einen HR-Experten für finale Entscheidungen.</small></p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
